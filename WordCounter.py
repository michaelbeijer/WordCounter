"""
Batch word counter with GUI for translators (FineCount-inspired).

Counts: .docx, .pptx, .xlsx, optional .pdf
With optional Apache Tika: 50+ additional formats (.doc, .odt, .html, .epub, .rtf, etc.)

Install:
  pip install python-docx python-pptx openpyxl
Optional PDF:
  pip install pdfminer.six
Optional Tika (50+ extra formats, requires Java):
  pip install tika
"""
from __future__ import annotations

import os
import sys
import re
import json
import threading
import queue
import csv
from datetime import datetime
from dataclasses import dataclass
from typing import List, Optional, Tuple, Dict

# --- Bundled JRE + Tika JAR detection (for PyInstaller frozen builds) ---
def _setup_bundled_tika():
    """Configure tika-python to use the JRE and Tika JAR bundled with the EXE."""
    if not getattr(sys, 'frozen', False):
        return
    # _MEIPASS points to _internal/ in onedir mode
    bundle_dir = sys._MEIPASS
    java_exe = os.path.join(bundle_dir, 'jre', 'bin', 'java.exe')
    tika_jar = os.path.join(bundle_dir, 'tika', 'tika-server-standard-3.1.0.jar')
    if os.path.isfile(java_exe):
        os.environ['TIKA_JAVA'] = java_exe
    if os.path.isfile(tika_jar):
        os.environ['TIKA_SERVER_JAR'] = 'file:///' + tika_jar.replace('\\', '/')

_setup_bundled_tika()

import tkinter as tk
from tkinter import ttk, filedialog, messagebox

# --- Optional imports ---
DOCX_OK = PPTX_OK = XLSX_OK = PDF_OK = False

try:
    from docx import Document
    DOCX_OK = True
except Exception:
    DOCX_OK = False

try:
    from pptx import Presentation
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
        PPTX_PLACEHOLDER_TYPES = PP_PLACEHOLDER_TYPE
    except Exception:
        PPTX_PLACEHOLDER_TYPES = None
    PPTX_OK = True
except Exception:
    PPTX_OK = False

try:
    import openpyxl
    XLSX_OK = True
except Exception:
    XLSX_OK = False

try:
    from pdfminer.high_level import extract_text as pdf_extract_text
    PDF_OK = True
except Exception:
    PDF_OK = False

TIKA_OK = False
try:
    import logging as _logging
    _logging.getLogger("tika").setLevel(_logging.WARNING)
    from tika import parser as tika_parser
    TIKA_OK = True
except Exception:
    TIKA_OK = False

APP_NAME = "WordCounter"
APP_AUTHOR = "Michael Beijer"
APP_VERSION = "0.6.0"


# ---------------- Tokenisation/stat helpers ----------------
WORD_RE = re.compile(r"[A-Za-zÀ-ÖØ-öø-ÿ0-9]+(?:['’][A-Za-zÀ-ÖØ-öø-ÿ0-9]+)?")
NUM_RE = re.compile(r"\b\d+(?:[.,]\d+)?\b")

def safe_join_text(parts: List[str]) -> str:
    return "\n".join(p for p in parts if p and p.strip())

def count_words(text: str) -> int:
    return len(WORD_RE.findall(text or ""))

def count_chars_with_spaces(text: str) -> int:
    return len(text or "")

def count_chars_no_spaces(text: str) -> int:
    return len(re.sub(r"\s+", "", text or ""))

def count_numbers(text: str) -> int:
    return len(NUM_RE.findall(text or ""))

def count_sentences(text: str) -> int:
    # Simple heuristic: count sentence end punctuation.
    t = (text or "").strip()
    if not t:
        return 0
    # Split on . ! ? followed by whitespace/end; avoid counting ellipses heavily.
    parts = re.split(r"(?<=[.!?])\s+", t)
    return sum(1 for p in parts if p.strip())

def count_paragraphs(text: str) -> int:
    # paragraphs separated by blank lines
    t = (text or "").strip()
    if not t:
        return 0
    blocks = re.split(r"\n\s*\n+", t)
    return sum(1 for b in blocks if b.strip())


# ---------------- Segment helpers (for repetition analysis) ----------------
SENTENCE_SPLIT_RE = re.compile(r'(?<=[.!?])\s+')

def normalize_segment(seg: str) -> str:
    """Normalize a segment for repetition comparison: collapse whitespace, strip."""
    return re.sub(r'\s+', ' ', seg.strip())

def text_to_sentences(text: str) -> List[str]:
    """Split text into sentence segments for repetition analysis (non-translation formats)."""
    if not text or not text.strip():
        return []
    sentences = SENTENCE_SPLIT_RE.split(text.strip())
    return [s.strip() for s in sentences if s.strip()]


# ---------------- Settings ----------------
@dataclass
class Settings:
    include_subfolders: bool = True

    # DOCX
    docx_include_body: bool = True
    docx_include_tables: bool = True
    docx_include_headers: bool = False
    docx_include_footers: bool = False

    # PPTX
    pptx_include_slide_text: bool = True
    pptx_include_footer_placeholders: bool = False
    pptx_include_speaker_notes: bool = False

    # XLSX
    xlsx_include_text: bool = True
    xlsx_include_numbers: bool = False
    xlsx_include_comments: bool = False
    xlsx_include_hidden_sheets: bool = False

    # PDF
    pdf_include: bool = True
    pdf_remove_repeating_headers_footers: bool = True

    # Translation files (SDLXLIFF, XLIFF, TMX, PO)
    xliff_count_target: bool = False  # False = count source, True = count target

    # Pages estimate
    words_per_page: int = 330  # common translation estimate; adjustable


# ---------------- DOCX ----------------
def _docx_collect(container, include_tables: bool) -> List[str]:
    parts: List[str] = []
    for p in getattr(container, "paragraphs", []):
        t = getattr(p, "text", "")
        if t and t.strip():
            parts.append(t)
    if include_tables:
        for table in getattr(container, "tables", []):
            for row in table.rows:
                for cell in row.cells:
                    ct = getattr(cell, "text", "")
                    if ct and ct.strip():
                        parts.append(ct)
    return parts

def docx_text(path: str, s: Settings) -> str:
    doc = Document(path)
    parts: List[str] = []

    if s.docx_include_body:
        parts.extend(_docx_collect(doc, include_tables=s.docx_include_tables))

    if s.docx_include_headers or s.docx_include_footers:
        for sec in doc.sections:
            if s.docx_include_headers:
                parts.extend(_docx_collect(sec.header, include_tables=s.docx_include_tables))
            if s.docx_include_footers:
                parts.extend(_docx_collect(sec.footer, include_tables=s.docx_include_tables))

    return safe_join_text(parts)

def extract_docx(path: str, s: Settings) -> ExtractionResult:
    if not DOCX_OK:
        return ExtractionResult("", "python-docx not installed")
    try:
        return ExtractionResult(docx_text(path, s), None)
    except Exception as e:
        return ExtractionResult("", f"DOCX error: {e}")


# ---------------- PPTX ----------------
def is_footer_placeholder(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
    except Exception:
        return False
    try:
        pht = shape.placeholder_format.type
        if PPTX_PLACEHOLDER_TYPES is not None:
            footer_types = {
                PPTX_PLACEHOLDER_TYPES.DATE_AND_TIME,
                PPTX_PLACEHOLDER_TYPES.FOOTER,
                PPTX_PLACEHOLDER_TYPES.SLIDE_NUMBER,
            }
            return pht in footer_types
        name = str(pht).upper()
        return any(k in name for k in ["DATE", "FOOTER", "SLIDE_NUMBER"])
    except Exception:
        return False

def pptx_text(path: str, s: Settings) -> str:
    prs = Presentation(path)
    parts: List[str] = []

    for slide in prs.slides:
        if s.pptx_include_slide_text:
            for shape in slide.shapes:
                if (not s.pptx_include_footer_placeholders) and is_footer_placeholder(shape):
                    continue
                try:
                    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                        t = shape.text
                        if t and t.strip():
                            parts.append(t)
                except Exception:
                    continue

        if s.pptx_include_speaker_notes:
            try:
                ns = slide.notes_slide
                if ns and ns.notes_text_frame:
                    nt = ns.notes_text_frame.text
                    if nt and nt.strip():
                        parts.append(nt)
            except Exception:
                pass

    return safe_join_text(parts)

def extract_pptx(path: str, s: Settings) -> ExtractionResult:
    if not PPTX_OK:
        return ExtractionResult("", "python-pptx not installed")
    try:
        return ExtractionResult(pptx_text(path, s), None)
    except Exception as e:
        return ExtractionResult("", f"PPTX error: {e}")


# ---------------- XLSX ----------------
def xlsx_text(path: str, s: Settings) -> str:
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts: List[str] = []

    for ws in wb.worksheets:
        if (not s.xlsx_include_hidden_sheets) and getattr(ws, "sheet_state", "visible") != "visible":
            continue

        for row in ws.iter_rows(values_only=False):
            for cell in row:
                v = cell.value

                if s.xlsx_include_text and isinstance(v, str) and v.strip():
                    parts.append(v)

                if s.xlsx_include_numbers and isinstance(v, (int, float)):
                    parts.append(str(v))

                if s.xlsx_include_comments:
                    try:
                        cmt = cell.comment
                        if cmt and cmt.text and cmt.text.strip():
                            parts.append(cmt.text)
                    except Exception:
                        pass

    return safe_join_text(parts)

def extract_xlsx(path: str, s: Settings) -> ExtractionResult:
    if not XLSX_OK:
        return ExtractionResult("", "openpyxl not installed")
    try:
        return ExtractionResult(xlsx_text(path, s), None)
    except Exception as e:
        return ExtractionResult("", f"XLSX error: {e}")


# ---------------- PDF ----------------
def _remove_repeating_lines(text: str) -> str:
    pages = text.split("\f")
    if len(pages) <= 1:
        return text

    line_freq: Dict[str, int] = {}
    page_norm_lines: List[List[str]] = []

    for p in pages:
        lines = [ln.strip() for ln in p.splitlines() if ln.strip()]
        norm = [re.sub(r"\d+", "#", ln) for ln in lines]
        page_norm_lines.append(norm)
        for ln in set(norm):
            line_freq[ln] = line_freq.get(ln, 0) + 1

    threshold = max(2, int(0.6 * len(pages)))
    repeating = {ln for ln, n in line_freq.items() if n >= threshold}

    cleaned_pages: List[str] = []
    for original_page in pages:
        cleaned_lines: List[str] = []
        for ln in original_page.splitlines():
            stripped = ln.strip()
            if not stripped:
                continue
            norm = re.sub(r"\d+", "#", stripped)
            if norm in repeating:
                continue
            cleaned_lines.append(stripped)
        cleaned_pages.append("\n".join(cleaned_lines))

    return "\n".join(cleaned_pages)

def extract_pdf(path: str, s: Settings) -> ExtractionResult:
    if not s.pdf_include:
        return ExtractionResult("", "PDF disabled")
    if not PDF_OK:
        return ExtractionResult("", "pdfminer.six not installed (PDF skipped)")
    try:
        text = pdf_extract_text(path) or ""
        if s.pdf_remove_repeating_headers_footers:
            text = _remove_repeating_lines(text)
        return ExtractionResult(text, None)
    except Exception as e:
        return ExtractionResult("", f"PDF error: {e}")


# ---------------- Translation formats (dedicated parsers) ----------------
import xml.etree.ElementTree as ET

# XLIFF inline tag names that are placeholders (no translatable text inside)
_XLIFF_SKIP_TAGS = {"bpt", "ept", "ph", "it", "x"}

def _local_tag(elem) -> str:
    """Return the local tag name (without namespace)."""
    tag = elem.tag
    if tag.startswith("{"):
        return tag[tag.index("}") + 1:]
    return tag

def _xml_itertext(elem) -> str:
    """Recursively extract translatable text, skipping XLIFF placeholder tags."""
    parts = []
    if elem.text:
        parts.append(elem.text)
    for child in elem:
        ltag = _local_tag(child)
        if ltag in _XLIFF_SKIP_TAGS:
            # Skip the tag content but keep the tail (text after the tag)
            if child.tail:
                parts.append(child.tail)
        else:
            # <g>, <mrk>, <sub>, etc. — recurse to get text inside
            parts.append(_xml_itertext(child))
            if child.tail:
                parts.append(child.tail)
    return "".join(parts)

def _detect_xliff_ns(root) -> str:
    """Detect the XLIFF namespace from the root element."""
    tag = root.tag
    if tag.startswith("{"):
        return tag[1:tag.index("}")]
    return ""

def _xliff_iter(root, ns: str, tag_name: str):
    """Iterate over elements with the given tag name, namespace-aware."""
    if ns:
        return root.iter(f"{{{ns}}}{tag_name}")
    return root.iter(tag_name)

def _xliff_findall(elem, ns: str, tag_name: str):
    """Find all children with the given tag name, namespace-aware."""
    if ns:
        return elem.findall(f"{{{ns}}}{tag_name}")
    return elem.findall(tag_name)

def _extract_mrk_segments(elem, ns: str) -> List[str]:
    """Extract text from <mrk mtype='seg'> children, or fallback to full element text."""
    mrks = _xliff_findall(elem, ns, "mrk")
    segments = []
    if mrks:
        for mrk in mrks:
            if mrk.get("mtype") == "seg":
                text = _xml_itertext(mrk).strip()
                if text:
                    segments.append(text)
    else:
        text = _xml_itertext(elem).strip()
        if text:
            segments.append(text)
    return segments

def extract_sdlxliff(path: str, count_target: bool = False) -> ExtractionResult:
    """Extract source or target segments from SDL Trados .sdlxliff files."""
    try:
        tree = ET.parse(path)
        root = tree.getroot()
        ns = _detect_xliff_ns(root)

        # Detect languages from <file> element
        file_elem = next(_xliff_iter(root, ns, "file"), None)
        src_lang = file_elem.get("source-language", "?") if file_elem is not None else "?"
        tgt_lang = file_elem.get("target-language", "?") if file_elem is not None else "?"
        counting = "target" if count_target else "source"
        lang_label = tgt_lang if count_target else src_lang

        segments = []
        if count_target:
            # Extract from <target> elements within <trans-unit>
            for tu in _xliff_iter(root, ns, "trans-unit"):
                target = tu.find(f"{{{ns}}}target" if ns else "target")
                if target is not None:
                    segments.extend(_extract_mrk_segments(target, ns))
        else:
            # Extract from <seg-source> (preferred) or <source>
            for seg_src in _xliff_iter(root, ns, "seg-source"):
                segments.extend(_extract_mrk_segments(seg_src, ns))
            if not segments:
                for src in _xliff_iter(root, ns, "source"):
                    segments.extend(_extract_mrk_segments(src, ns))

        if not segments:
            return ExtractionResult("", f"No {counting} segments found in SDLXLIFF")
        note = f"SDLXLIFF {counting} [{lang_label}]"
        return ExtractionResult("\n".join(segments), note, segments)
    except Exception as e:
        return ExtractionResult("", f"SDLXLIFF error: {e}")

def extract_xliff(path: str, count_target: bool = False) -> ExtractionResult:
    """Extract source or target segments from XLIFF (.xliff, .xlf, .mqxliff) files."""
    try:
        tree = ET.parse(path)
        root = tree.getroot()
        ns = _detect_xliff_ns(root)
        counting = "target" if count_target else "source"

        # Detect languages
        file_elem = next(_xliff_iter(root, ns, "file"), None)
        src_lang = file_elem.get("source-language", "?") if file_elem is not None else "?"
        tgt_lang = file_elem.get("target-language", "?") if file_elem is not None else "?"
        lang_label = tgt_lang if count_target else src_lang

        segments = []
        if count_target:
            for tu in _xliff_iter(root, ns, "trans-unit"):
                target = tu.find(f"{{{ns}}}target" if ns else "target")
                if target is not None:
                    segments.extend(_extract_mrk_segments(target, ns))
        else:
            # Try <seg-source> first, fall back to <source>
            for seg_src in _xliff_iter(root, ns, "seg-source"):
                segments.extend(_extract_mrk_segments(seg_src, ns))
            if not segments:
                for src in _xliff_iter(root, ns, "source"):
                    segments.extend(_extract_mrk_segments(src, ns))

        if not segments:
            return ExtractionResult("", f"No {counting} segments found in XLIFF")
        note = f"XLIFF {counting} [{lang_label}]"
        return ExtractionResult("\n".join(segments), note, segments)
    except Exception as e:
        return ExtractionResult("", f"XLIFF error: {e}")

def extract_tmx(path: str) -> ExtractionResult:
    """Extract source segments from TMX files (first language variant per TU)."""
    try:
        tree = ET.parse(path)
        root = tree.getroot()
        # Detect source language from header
        header = root.find(".//header")
        srclang = header.get("srclang", "") if header is not None else ""

        segments = []
        for tu in root.iter("tu"):
            tuvs = tu.findall("tuv")
            source_tuv = None
            if srclang:
                # Find TUV matching source language
                for tuv in tuvs:
                    lang = tuv.get("{http://www.w3.org/XML/1998/namespace}lang", "") or tuv.get("lang", "")
                    if lang.lower().startswith(srclang.lower()):
                        source_tuv = tuv
                        break
            if source_tuv is None and tuvs:
                source_tuv = tuvs[0]  # First TUV = source
            if source_tuv is not None:
                seg = source_tuv.find("seg")
                if seg is not None:
                    text = _xml_itertext(seg).strip()
                    if text:
                        segments.append(text)

        if not segments:
            return ExtractionResult("", "No source segments found in TMX")
        return ExtractionResult("\n".join(segments), None, segments)
    except Exception as e:
        return ExtractionResult("", f"TMX error: {e}")

def extract_po(path: str) -> ExtractionResult:
    """Extract source strings (msgid) from PO/POT files."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        segments = []
        # Match msgid entries (possibly multi-line)
        in_msgid = False
        current = []
        for line in content.splitlines():
            stripped = line.strip()
            if stripped.startswith("msgid "):
                in_msgid = True
                # Extract the string after msgid
                val = stripped[6:].strip().strip('"')
                if val:
                    current.append(val)
            elif in_msgid and stripped.startswith('"') and stripped.endswith('"'):
                current.append(stripped[1:-1])
            else:
                if in_msgid and current:
                    text = "".join(current)
                    if text:  # Skip empty msgid ""
                        segments.append(text)
                    current = []
                in_msgid = False
        # Flush last
        if in_msgid and current:
            text = "".join(current)
            if text:
                segments.append(text)

        if not segments:
            return ExtractionResult("", "No source strings found in PO file")
        return ExtractionResult("\n".join(segments), None, segments)
    except Exception as e:
        return ExtractionResult("", f"PO error: {e}")

# Translation format extensions handled by dedicated parsers
TRANSLATION_EXTS = {".sdlxliff", ".xliff", ".xlf", ".mqxliff", ".tmx", ".po", ".pot"}


# ---------------- Tika (universal fallback) ----------------
TIKA_EXTS = {
    # Legacy Microsoft Office
    ".doc", ".xls", ".ppt",
    # Rich Text Format
    ".rtf",
    # OpenDocument
    ".odt", ".odp", ".ods", ".odg",
    # Web / markup
    ".html", ".htm", ".xhtml", ".xml",
    # Plain text
    ".txt", ".csv", ".tsv",
    # Markup text
    ".md", ".rst", ".tex", ".latex",
    # E-books
    ".epub",
    # Email
    ".eml", ".msg",
    # Translation / localisation
    ".xliff", ".xlf", ".tmx", ".sdlxliff", ".mqxliff",
    ".po", ".pot", ".tbx",
    # Subtitles
    ".srt", ".vtt", ".ass", ".sub",
    # Desktop publishing
    ".idml",
    # Visio
    ".vsdx",
    # Images (OCR — requires Tesseract on the system)
    ".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".gif",
    # Data formats
    ".json", ".yaml", ".yml",
    # Localisation / config
    ".properties", ".strings", ".resx",
}

def extract_tika(path: str) -> ExtractionResult:
    if not TIKA_OK:
        return ExtractionResult("", "Apache Tika not installed")
    try:
        parsed = tika_parser.from_file(path)
        text = parsed.get("content") or ""
        return ExtractionResult(text.strip(), None)
    except Exception as e:
        return ExtractionResult("", f"Tika error: {e}")


# ---------------- Batch + metrics ----------------
CORE_EXTS = {".docx", ".pptx", ".xlsx", ".pdf"} | TRANSLATION_EXTS

def get_supported_exts(include_pdfs: bool = True) -> set:
    exts = set(CORE_EXTS)
    if not include_pdfs:
        exts.discard(".pdf")
    if TIKA_OK:
        exts |= TIKA_EXTS
    return exts

@dataclass
class ExtractionResult:
    text: str
    note: Optional[str]
    segments: Optional[List[str]] = None  # individual segments for repetition analysis

@dataclass
class FileMetrics:
    filepath: str
    words: int
    chars: int
    chars_nospace: int
    numbers: int
    sentences: int
    paragraphs: int
    pages_est: float
    note: str = ""
    text: str = ""
    segments: Optional[List[str]] = None  # segments from extraction

@dataclass
class RepetitionInfo:
    """Per-file repetition breakdown."""
    total_segments: int
    unique_segments: int
    repeated_segments: int
    unique_words: int
    repeated_words: int
    unique_chars: int
    repeated_chars: int

@dataclass
class BatchRepetitionResult:
    """Corpus-level repetition analysis results."""
    per_file: Dict[str, RepetitionInfo]
    corpus_unique_segments: int
    corpus_total_segments: int
    corpus_repeated_segments: int
    corpus_unique_words: int
    corpus_repeated_words: int
    corpus_unique_chars: int
    corpus_repeated_chars: int

def compute_metrics(text: str, s: Settings) -> Tuple[int, int, int, int, int, int, float]:
    w = count_words(text)
    c = count_chars_with_spaces(text)
    cns = count_chars_no_spaces(text)
    n = count_numbers(text)
    sent = count_sentences(text)
    para = count_paragraphs(text)
    pages = (w / s.words_per_page) if s.words_per_page > 0 else 0.0
    return w, c, cns, n, sent, para, pages

def analyze_repetitions(results: List[FileMetrics]) -> BatchRepetitionResult:
    """Analyze cross-document segment repetitions across all files in a batch.

    First occurrence of each segment = unique, subsequent = repetition.
    File processing order determines which file 'owns' unique segments.
    """
    seen: Dict[str, bool] = {}  # normalized_segment -> True (just presence)
    per_file: Dict[str, RepetitionInfo] = {}

    corpus_unique_words = 0
    corpus_repeated_words = 0
    corpus_unique_chars = 0
    corpus_repeated_chars = 0
    corpus_total_segments = 0
    corpus_repeated_segments = 0

    for fm in results:
        # Use stored segments for translation formats, sentence-split for others
        segments = fm.segments if fm.segments is not None else text_to_sentences(fm.text)

        file_unique_segs = 0
        file_repeated_segs = 0
        file_unique_words = 0
        file_repeated_words = 0
        file_unique_chars = 0
        file_repeated_chars = 0

        for seg in segments:
            norm = normalize_segment(seg)
            if not norm:
                continue

            seg_words = count_words(seg)
            seg_chars = count_chars_with_spaces(seg)
            corpus_total_segments += 1

            if norm in seen:
                file_repeated_segs += 1
                file_repeated_words += seg_words
                file_repeated_chars += seg_chars
                corpus_repeated_segments += 1
                corpus_repeated_words += seg_words
                corpus_repeated_chars += seg_chars
            else:
                file_unique_segs += 1
                file_unique_words += seg_words
                file_unique_chars += seg_chars
                corpus_unique_words += seg_words
                corpus_unique_chars += seg_chars
                seen[norm] = True

        per_file[fm.filepath] = RepetitionInfo(
            total_segments=file_unique_segs + file_repeated_segs,
            unique_segments=file_unique_segs,
            repeated_segments=file_repeated_segs,
            unique_words=file_unique_words,
            repeated_words=file_repeated_words,
            unique_chars=file_unique_chars,
            repeated_chars=file_repeated_chars,
        )

    return BatchRepetitionResult(
        per_file=per_file,
        corpus_unique_segments=len(seen),
        corpus_total_segments=corpus_total_segments,
        corpus_repeated_segments=corpus_repeated_segments,
        corpus_unique_words=corpus_unique_words,
        corpus_repeated_words=corpus_repeated_words,
        corpus_unique_chars=corpus_unique_chars,
        corpus_repeated_chars=corpus_repeated_chars,
    )

def extract_text_by_type(path: str, s: Settings) -> ExtractionResult:
    ext = os.path.splitext(path)[1].lower()
    # Dedicated extractors (with fine-grained settings)
    if ext == ".docx":
        return extract_docx(path, s)
    if ext == ".pptx":
        return extract_pptx(path, s)
    if ext == ".xlsx":
        return extract_xlsx(path, s)
    if ext == ".pdf":
        return extract_pdf(path, s)
    # Translation formats (dedicated XML/text parsers — always preferred over Tika)
    if ext == ".sdlxliff":
        return extract_sdlxliff(path, count_target=s.xliff_count_target)
    if ext in {".xliff", ".xlf", ".mqxliff"}:
        return extract_xliff(path, count_target=s.xliff_count_target)
    if ext == ".tmx":
        return extract_tmx(path)
    if ext in {".po", ".pot"}:
        return extract_po(path)
    # Tika fallback for all other formats
    if TIKA_OK:
        return extract_tika(path)
    return ExtractionResult("", "Unsupported file type")

def iter_files(folder: str, include_subfolders: bool, include_pdfs: bool) -> List[str]:
    supported = get_supported_exts(include_pdfs)
    def ok_ext(fn: str) -> bool:
        return os.path.splitext(fn)[1].lower() in supported

    paths: List[str] = []
    if include_subfolders:
        for root, _, files in os.walk(folder):
            for fn in files:
                if ok_ext(fn):
                    paths.append(os.path.join(root, fn))
    else:
        for fn in os.listdir(folder):
            p = os.path.join(folder, fn)
            if os.path.isfile(p) and ok_ext(fn):
                paths.append(p)
    return sorted(paths)

def filter_supported(files: List[str], include_pdfs: bool) -> List[str]:
    supported = get_supported_exts(include_pdfs)
    out = []
    for p in files:
        if os.path.splitext(p)[1].lower() in supported:
            out.append(p)
    return sorted(list(dict.fromkeys(out)))  # de-dupe preserve order


# ---------------- GUI ----------------
class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title(f"{APP_NAME} v{APP_VERSION}, by {APP_AUTHOR}")
        self.geometry("1200x760")

        self.folder_var = tk.StringVar(value="")
        self._queue = queue.Queue()
        self._results: List[FileMetrics] = []
        self._file_list: List[str] = []
        self._repetition: Optional[BatchRepetitionResult] = None

        # --- Settings vars (counting) ---
        self.include_subfolders_var = tk.BooleanVar(value=True)

        self.docx_body_var = tk.BooleanVar(value=True)
        self.docx_tables_var = tk.BooleanVar(value=True)
        self.docx_headers_var = tk.BooleanVar(value=False)
        self.docx_footers_var = tk.BooleanVar(value=False)

        self.pptx_slide_text_var = tk.BooleanVar(value=True)
        self.pptx_footer_ph_var = tk.BooleanVar(value=False)
        self.pptx_notes_var = tk.BooleanVar(value=False)

        self.xlsx_text_var = tk.BooleanVar(value=True)
        self.xlsx_numbers_var = tk.BooleanVar(value=False)
        self.xlsx_comments_var = tk.BooleanVar(value=False)
        self.xlsx_hidden_sheets_var = tk.BooleanVar(value=False)

        self.pdf_include_var = tk.BooleanVar(value=True)
        self.pdf_strip_repeat_var = tk.BooleanVar(value=True)

        self.xliff_count_target_var = tk.BooleanVar(value=False)

        self.words_per_page_var = tk.IntVar(value=330)

        # --- Billing vars ---
        self.bill_by_var = tk.StringVar(value="Words")  # Words / Characters / Pages (est.)
        self.rate_var = tk.DoubleVar(value=0.0)
        self.rep_rate_var = tk.DoubleVar(value=0.0)  # rate for repeated segments (0 = excluded)
        self.currency_var = tk.StringVar(value="GBP")
        self.tax_var = tk.DoubleVar(value=0.0)       # percent
        self.discount_var = tk.DoubleVar(value=0.0)  # percent

        self._load_settings()
        self._build_ui()
        self._poll_queue()
        self.protocol("WM_DELETE_WINDOW", self._on_close)

        # Auto-save settings when any variable changes
        for v in (
            self.include_subfolders_var,
            self.docx_body_var, self.docx_tables_var, self.docx_headers_var, self.docx_footers_var,
            self.pptx_slide_text_var, self.pptx_footer_ph_var, self.pptx_notes_var,
            self.xlsx_text_var, self.xlsx_numbers_var, self.xlsx_comments_var, self.xlsx_hidden_sheets_var,
            self.pdf_include_var, self.pdf_strip_repeat_var,
            self.xliff_count_target_var,
            self.words_per_page_var,
            self.bill_by_var, self.rate_var, self.rep_rate_var, self.currency_var, self.tax_var, self.discount_var,
        ):
            v.trace_add("write", lambda *_: self._save_settings())

    # -------- Settings persistence --------
    @staticmethod
    def _settings_file() -> str:
        return os.path.join(os.path.expanduser("~"), ".wordcounter_settings.json")

    def _save_settings(self):
        data = {
            "include_subfolders": self.include_subfolders_var.get(),
            "docx_body": self.docx_body_var.get(),
            "docx_tables": self.docx_tables_var.get(),
            "docx_headers": self.docx_headers_var.get(),
            "docx_footers": self.docx_footers_var.get(),
            "pptx_slide_text": self.pptx_slide_text_var.get(),
            "pptx_footer_ph": self.pptx_footer_ph_var.get(),
            "pptx_notes": self.pptx_notes_var.get(),
            "xlsx_text": self.xlsx_text_var.get(),
            "xlsx_numbers": self.xlsx_numbers_var.get(),
            "xlsx_comments": self.xlsx_comments_var.get(),
            "xlsx_hidden_sheets": self.xlsx_hidden_sheets_var.get(),
            "pdf_include": self.pdf_include_var.get(),
            "pdf_strip_repeat": self.pdf_strip_repeat_var.get(),
            "xliff_count_target": self.xliff_count_target_var.get(),
            "words_per_page": self.words_per_page_var.get(),
            "bill_by": self.bill_by_var.get(),
            "rate": self.rate_var.get(),
            "rep_rate": self.rep_rate_var.get(),
            "currency": self.currency_var.get(),
            "tax": self.tax_var.get(),
            "discount": self.discount_var.get(),
            "last_folder": self.folder_var.get(),
            "geometry": self.geometry(),
        }
        try:
            with open(self._settings_file(), "w", encoding="utf-8") as f:
                json.dump(data, f, indent=2)
        except Exception:
            pass  # non-critical

    def _load_settings(self):
        try:
            with open(self._settings_file(), "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception:
            return

        def _b(key, var):
            if key in data:
                var.set(bool(data[key]))

        def _s(key, var):
            if key in data:
                var.set(str(data[key]))

        def _n(key, var):
            if key in data:
                try:
                    var.set(data[key])
                except Exception:
                    pass

        _b("include_subfolders", self.include_subfolders_var)
        _b("docx_body", self.docx_body_var)
        _b("docx_tables", self.docx_tables_var)
        _b("docx_headers", self.docx_headers_var)
        _b("docx_footers", self.docx_footers_var)
        _b("pptx_slide_text", self.pptx_slide_text_var)
        _b("pptx_footer_ph", self.pptx_footer_ph_var)
        _b("pptx_notes", self.pptx_notes_var)
        _b("xlsx_text", self.xlsx_text_var)
        _b("xlsx_numbers", self.xlsx_numbers_var)
        _b("xlsx_comments", self.xlsx_comments_var)
        _b("xlsx_hidden_sheets", self.xlsx_hidden_sheets_var)
        _b("pdf_include", self.pdf_include_var)
        _b("pdf_strip_repeat", self.pdf_strip_repeat_var)
        _b("xliff_count_target", self.xliff_count_target_var)
        _n("words_per_page", self.words_per_page_var)
        _s("bill_by", self.bill_by_var)
        _n("rate", self.rate_var)
        _n("rep_rate", self.rep_rate_var)
        _s("currency", self.currency_var)
        _n("tax", self.tax_var)
        _n("discount", self.discount_var)
        _s("last_folder", self.folder_var)

        if "geometry" in data:
            try:
                self.geometry(data["geometry"])
            except Exception:
                pass

    def show_about(self):
        dlg = tk.Toplevel(self)
        dlg.title(f"About {APP_NAME}")
        dlg.resizable(False, False)
        dlg.grab_set()
        dlg.geometry(f"+{self.winfo_rootx() + 200}+{self.winfo_rooty() + 120}")

        frame = ttk.Frame(dlg, padding=24)
        frame.pack()

        ttk.Label(frame, text=f"{APP_NAME} v{APP_VERSION}",
                  font=("Segoe UI", 14, "bold")).pack(pady=(0, 8))
        ttk.Label(frame, text=f"by {APP_AUTHOR}").pack()

        # Clickable website link
        website_label = tk.Label(frame, text="michaelbeijer.co.uk",
                                 fg="blue", cursor="hand2", font=("Segoe UI", 9, "underline"))
        website_label.pack(pady=(2, 8))
        website_label.bind("<Button-1>", lambda e: __import__("webbrowser").open("https://michaelbeijer.co.uk/"))

        ttk.Separator(frame).pack(fill="x", pady=8)

        ttk.Label(frame, text="A batch word counter for translators.",
                  wraplength=300).pack(pady=(0, 8))

        # Clickable repo link
        repo_label = tk.Label(frame, text="GitHub: michaelbeijer/WordCounter",
                              fg="blue", cursor="hand2", font=("Segoe UI", 9, "underline"))
        repo_label.pack(pady=(0, 12))
        repo_label.bind("<Button-1>", lambda e: __import__("webbrowser").open("https://github.com/michaelbeijer/WordCounter"))

        ttk.Button(frame, text="Close", command=dlg.destroy).pack()

    def _on_close(self):
        self._save_settings()
        self.destroy()

    def _dependency_status(self) -> str:
        bits = [
            f"DOCX: {'OK' if DOCX_OK else 'missing'}",
            f"PPTX: {'OK' if PPTX_OK else 'missing'}",
            f"XLSX: {'OK' if XLSX_OK else 'missing'}",
            f"PDF: {'OK' if PDF_OK else 'missing (optional)'}",
            f"Tika: {'OK (+{0} formats)'.format(len(TIKA_EXTS)) if TIKA_OK else 'not installed (optional)'}",
        ]
        return f"v{APP_VERSION} • " + " • ".join(bits)

    def _get_settings(self) -> Settings:
        return Settings(
            include_subfolders=self.include_subfolders_var.get(),
            docx_include_body=self.docx_body_var.get(),
            docx_include_tables=self.docx_tables_var.get(),
            docx_include_headers=self.docx_headers_var.get(),
            docx_include_footers=self.docx_footers_var.get(),
            pptx_include_slide_text=self.pptx_slide_text_var.get(),
            pptx_include_footer_placeholders=self.pptx_footer_ph_var.get(),
            pptx_include_speaker_notes=self.pptx_notes_var.get(),
            xlsx_include_text=self.xlsx_text_var.get(),
            xlsx_include_numbers=self.xlsx_numbers_var.get(),
            xlsx_include_comments=self.xlsx_comments_var.get(),
            xlsx_include_hidden_sheets=self.xlsx_hidden_sheets_var.get(),
            pdf_include=self.pdf_include_var.get(),
            pdf_remove_repeating_headers_footers=self.pdf_strip_repeat_var.get(),
            xliff_count_target=self.xliff_count_target_var.get(),
            words_per_page=int(self.words_per_page_var.get() or 330),
        )

    def _build_ui(self):
        # Top bar
        top = ttk.Frame(self, padding=10)
        top.pack(fill="x")

        ttk.Label(top, text="Folder:").grid(row=0, column=0, sticky="w")
        ttk.Entry(top, textvariable=self.folder_var).grid(row=0, column=1, sticky="we", padx=(8, 8))
        ttk.Button(top, text="Browse…", command=self.browse).grid(row=0, column=2)

        ttk.Button(top, text="Add files…", command=self.add_files).grid(row=1, column=0, sticky="w", pady=(6, 0))
        ttk.Button(top, text="Remove selected", command=self.remove_selected).grid(row=1, column=1, sticky="w", padx=(8, 0), pady=(6, 0))
        ttk.Button(top, text="Remove all", command=self.remove_all).grid(row=1, column=2, sticky="w", padx=(8, 0), pady=(6, 0))

        self.run_btn = ttk.Button(top, text="Count", command=self.run_count)
        self.run_btn.grid(row=0, column=3, padx=(10, 0))

        self.export_btn = ttk.Button(top, text="Export CSV…", command=self.export_csv, state="disabled")
        self.export_btn.grid(row=0, column=4, padx=(8, 0))

        self.export_md_btn = ttk.Button(top, text="Export MD…", command=self.export_md, state="disabled")
        self.export_md_btn.grid(row=0, column=5, padx=(8, 0))

        self.copy_btn = ttk.Button(top, text="Copy report", command=self.copy_report, state="disabled")
        self.copy_btn.grid(row=0, column=6, padx=(8, 0))

        ttk.Button(top, text="About", command=self.show_about).grid(row=0, column=7, padx=(8, 0))

        top.columnconfigure(1, weight=1)

        ttk.Checkbutton(top, text="Include subfolders (for folder counts)", variable=self.include_subfolders_var)\
            .grid(row=1, column=4, columnspan=3, sticky="w", padx=(12, 0), pady=(6, 0))

        # Settings + Billing panel
        mid = ttk.Frame(self, padding=(10, 0, 10, 6))
        mid.pack(fill="x")

        settings = ttk.LabelFrame(mid, text="Counting settings", padding=10)
        settings.pack(side="left", fill="x", expand=True)

        billing = ttk.LabelFrame(mid, text="Billing", padding=10)
        billing.pack(side="right", fill="y", padx=(10, 0))

        # Settings grid
        # DOCX
        docx = ttk.LabelFrame(settings, text="Word (.docx)", padding=8)
        docx.grid(row=0, column=0, sticky="nwe", padx=(0, 8))
        ttk.Checkbutton(docx, text="Body", variable=self.docx_body_var).grid(row=0, column=0, sticky="w")
        ttk.Checkbutton(docx, text="Tables", variable=self.docx_tables_var).grid(row=1, column=0, sticky="w")
        ttk.Checkbutton(docx, text="Headers", variable=self.docx_headers_var).grid(row=2, column=0, sticky="w")
        ttk.Checkbutton(docx, text="Footers", variable=self.docx_footers_var).grid(row=3, column=0, sticky="w")

        # PPTX
        pptx = ttk.LabelFrame(settings, text="PowerPoint (.pptx)", padding=8)
        pptx.grid(row=0, column=1, sticky="nwe", padx=(0, 8))
        ttk.Checkbutton(pptx, text="Slide text", variable=self.pptx_slide_text_var).grid(row=0, column=0, sticky="w")
        ttk.Checkbutton(pptx, text="Footer/date/slide # placeholders", variable=self.pptx_footer_ph_var)\
            .grid(row=1, column=0, sticky="w")
        ttk.Checkbutton(pptx, text="Speaker notes", variable=self.pptx_notes_var).grid(row=2, column=0, sticky="w")

        # XLSX
        xlsx = ttk.LabelFrame(settings, text="Excel (.xlsx)", padding=8)
        xlsx.grid(row=0, column=2, sticky="nwe")
        ttk.Checkbutton(xlsx, text="Cell text", variable=self.xlsx_text_var).grid(row=0, column=0, sticky="w")
        ttk.Checkbutton(xlsx, text="Numbers", variable=self.xlsx_numbers_var).grid(row=1, column=0, sticky="w")
        ttk.Checkbutton(xlsx, text="Comments/notes", variable=self.xlsx_comments_var).grid(row=2, column=0, sticky="w")
        ttk.Checkbutton(xlsx, text="Hidden sheets", variable=self.xlsx_hidden_sheets_var).grid(row=3, column=0, sticky="w")

        # PDF + pages estimate
        bottom_settings = ttk.Frame(settings)
        bottom_settings.grid(row=1, column=0, columnspan=3, sticky="we", pady=(8, 0))
        ttk.Checkbutton(bottom_settings, text="Include PDFs", variable=self.pdf_include_var).grid(row=0, column=0, sticky="w")
        ttk.Checkbutton(bottom_settings, text="Remove repeating header/footer lines (PDF heuristic)", variable=self.pdf_strip_repeat_var)\
            .grid(row=0, column=1, sticky="w", padx=(10, 0))
        ttk.Label(bottom_settings, text="Words per page (estimate):").grid(row=0, column=2, sticky="e", padx=(18, 6))
        ttk.Spinbox(bottom_settings, from_=100, to=1000, textvariable=self.words_per_page_var, width=6)\
            .grid(row=0, column=3, sticky="w")
        ttk.Checkbutton(bottom_settings, text="XLIFF/SDLXLIFF: count target segments (default: source)",
                         variable=self.xliff_count_target_var).grid(row=1, column=0, columnspan=3, sticky="w", pady=(4, 0))

        settings.columnconfigure(0, weight=1)
        settings.columnconfigure(1, weight=1)
        settings.columnconfigure(2, weight=1)

        # Billing controls
        ttk.Label(billing, text="Bill by:").grid(row=0, column=0, sticky="w")
        bill_by = ttk.Combobox(billing, textvariable=self.bill_by_var, values=["Words", "Characters", "Pages (est.)"], state="readonly", width=14)
        bill_by.grid(row=0, column=1, sticky="w")

        ttk.Label(billing, text="Rate:").grid(row=1, column=0, sticky="w", pady=(6, 0))
        ttk.Entry(billing, textvariable=self.rate_var, width=10).grid(row=1, column=1, sticky="w", pady=(6, 0))

        ttk.Label(billing, text="Rep. rate:").grid(row=2, column=0, sticky="w", pady=(6, 0))
        ttk.Entry(billing, textvariable=self.rep_rate_var, width=10).grid(row=2, column=1, sticky="w", pady=(6, 0))

        ttk.Label(billing, text="Currency:").grid(row=3, column=0, sticky="w", pady=(6, 0))
        ttk.Entry(billing, textvariable=self.currency_var, width=10).grid(row=3, column=1, sticky="w", pady=(6, 0))

        ttk.Label(billing, text="Tax %:").grid(row=4, column=0, sticky="w", pady=(6, 0))
        ttk.Entry(billing, textvariable=self.tax_var, width=10).grid(row=4, column=1, sticky="w", pady=(6, 0))

        ttk.Label(billing, text="Discount %:").grid(row=5, column=0, sticky="w", pady=(6, 0))
        ttk.Entry(billing, textvariable=self.discount_var, width=10).grid(row=5, column=1, sticky="w", pady=(6, 0))

        ttk.Separator(billing).grid(row=6, column=0, columnspan=2, sticky="we", pady=10)

        self.total_billable_var = tk.StringVar(value="Unique units: 0")
        self.total_rep_billable_var = tk.StringVar(value="Rep. units: 0")
        self.total_amount_var = tk.StringVar(value="Total: 0.00")
        ttk.Label(billing, textvariable=self.total_billable_var).grid(row=7, column=0, columnspan=2, sticky="w")
        ttk.Label(billing, textvariable=self.total_rep_billable_var).grid(row=8, column=0, columnspan=2, sticky="w", pady=(2, 0))
        ttk.Label(billing, textvariable=self.total_amount_var, font=("Segoe UI", 11, "bold")).grid(row=9, column=0, columnspan=2, sticky="w", pady=(4, 0))

        # Dependency/status line
        self.status_var = tk.StringVar(value=self._dependency_status())
        ttk.Label(self, textvariable=self.status_var, padding=(10, 0)).pack(fill="x")

        # Results table
        frame = ttk.Frame(self, padding=10)
        frame.pack(fill="both", expand=True)

        cols = ("words", "uniq_w", "rep_w", "rep_pct", "chars", "chars_ns", "nums", "nums_pct", "sent", "para", "pages", "note", "path")
        self.tree = ttk.Treeview(frame, columns=cols, show="headings")

        headings = {
            "words": "Words",
            "uniq_w": "Unique",
            "rep_w": "Repeated",
            "rep_pct": "% rep.",
            "chars": "Chars",
            "chars_ns": "Chars (no sp)",
            "nums": "Numbers",
            "nums_pct": "% nums",
            "sent": "Sent.",
            "para": "Para.",
            "pages": "Pages (est.)",
            "note": "Note",
            "path": "File",
        }
        for k, h in headings.items():
            self.tree.heading(k, text=h)

        self.tree.column("words", width=80, anchor="e")
        self.tree.column("uniq_w", width=80, anchor="e")
        self.tree.column("rep_w", width=80, anchor="e")
        self.tree.column("rep_pct", width=60, anchor="e")
        self.tree.column("chars", width=90, anchor="e")
        self.tree.column("chars_ns", width=110, anchor="e")
        self.tree.column("nums", width=80, anchor="e")
        self.tree.column("nums_pct", width=70, anchor="e")
        self.tree.column("sent", width=60, anchor="e")
        self.tree.column("para", width=60, anchor="e")
        self.tree.column("pages", width=90, anchor="e")
        self.tree.column("note", width=220, anchor="w")
        self.tree.column("path", width=520, anchor="w")

        vsb = ttk.Scrollbar(frame, orient="vertical", command=self.tree.yview)
        self.tree.configure(yscrollcommand=vsb.set)

        self.tree.pack(side="left", fill="both", expand=True)
        vsb.pack(side="right", fill="y")

        # Totals bar
        bottom = ttk.Frame(self, padding=10)
        bottom.pack(fill="x")
        self.total_files_var = tk.StringVar(value="Files: 0")
        self.total_words_var = tk.StringVar(value="Total words: 0")
        self.total_unique_words_var = tk.StringVar(value="")
        self.total_rep_words_var = tk.StringVar(value="")
        ttk.Label(bottom, textvariable=self.total_files_var).pack(side="left")
        ttk.Label(bottom, textvariable=self.total_words_var).pack(side="left", padx=(16, 0))
        ttk.Label(bottom, textvariable=self.total_unique_words_var).pack(side="left", padx=(16, 0))
        ttk.Label(bottom, textvariable=self.total_rep_words_var).pack(side="left", padx=(16, 0))

        # Recompute billing when these change
        for v in (self.bill_by_var, self.rate_var, self.rep_rate_var, self.currency_var, self.tax_var, self.discount_var, self.words_per_page_var):
            v.trace_add("write", lambda *_: self.update_billing())

        # Keyboard shortcuts
        self.bind("<Control-o>", lambda e: self.browse())
        self.bind("<Control-Return>", lambda e: self.run_count())
        self.bind("<F5>", lambda e: self.run_count())

    # -------- File list management --------
    def _pick_files_dialog(self):
        if TIKA_OK:
            all_exts = " ".join(f"*{e}" for e in sorted(get_supported_exts()))
            types = [
                ("All supported", all_exts),
                ("Microsoft Office", "*.docx *.doc *.pptx *.ppt *.xlsx *.xls"),
                ("OpenDocument", "*.odt *.odp *.ods"),
                ("PDF", "*.pdf"),
                ("Text / Markup", "*.txt *.html *.htm *.xml *.md *.rst *.csv *.rtf"),
                ("Translation / L10n", "*.xliff *.xlf *.tmx *.sdlxliff *.mqxliff *.po *.pot *.tbx"),
                ("Subtitles", "*.srt *.vtt *.ass *.sub"),
                ("E-books", "*.epub"),
                ("Images (OCR)", "*.png *.jpg *.jpeg *.tiff *.tif *.bmp *.gif"),
                ("All files", "*.*"),
            ]
        else:
            types = [
                ("Supported", "*.docx *.pptx *.xlsx *.pdf"),
                ("Word", "*.docx"),
                ("PowerPoint", "*.pptx"),
                ("Excel", "*.xlsx"),
                ("PDF", "*.pdf"),
                ("All files", "*.*"),
            ]
        return filedialog.askopenfilenames(title="Select files", filetypes=types)

    def browse(self):
        dlg = tk.Toplevel(self)
        dlg.title("Browse")
        dlg.resizable(False, False)
        dlg.grab_set()
        dlg.geometry(f"+{self.winfo_rootx() + 150}+{self.winfo_rooty() + 80}")

        frame = ttk.Frame(dlg, padding=20)
        frame.pack()

        ttk.Label(frame, text="What would you like to select?",
                  font=("Segoe UI", 10)).pack(pady=(0, 14))

        def pick_files():
            dlg.destroy()
            s = self._get_settings()
            files = self._pick_files_dialog()
            if files:
                self._file_list = filter_supported(list(files), include_pdfs=s.pdf_include)
                if self._file_list:
                    first_dir = os.path.dirname(self._file_list[0])
                    self.folder_var.set(first_dir)
                self.status_var.set(f"Selected {len(self._file_list)} file(s). Click Count.")

        def pick_folder():
            dlg.destroy()
            folder = filedialog.askdirectory(title="Choose a folder")
            if folder:
                self.folder_var.set(folder)
                self._file_list = []
                self.status_var.set("Folder selected. Click Count.")

        btn_frame = ttk.Frame(frame)
        btn_frame.pack(fill="x")
        ttk.Button(btn_frame, text="Select Files…", command=pick_files).pack(fill="x", pady=(0, 6))
        ttk.Button(btn_frame, text="Select Folder…", command=pick_folder).pack(fill="x")

        dlg.wait_window()

    def add_folder(self):
        folder = self.folder_var.get().strip()
        if not folder or not os.path.isdir(folder):
            messagebox.showerror("Folder required", "Please select a valid folder first (Browse…).")
            return
        s = self._get_settings()
        paths = iter_files(folder, s.include_subfolders, include_pdfs=s.pdf_include)
        self._file_list = filter_supported(self._file_list + paths, include_pdfs=s.pdf_include)
        self.run_count()

    def add_files(self):
        s = self._get_settings()
        files = self._pick_files_dialog()
        if files:
            self._file_list = filter_supported(self._file_list + list(files), include_pdfs=s.pdf_include)
            self.run_count()

    def remove_selected(self):
        sel = self.tree.selection()
        if not sel:
            return
        paths_to_remove = set()
        for item in sel:
            paths_to_remove.add(self.tree.set(item, "path"))
        self._file_list = [p for p in self._file_list if p not in paths_to_remove]
        self.run_count()

    def remove_all(self):
        self._file_list = []
        self.clear_results()

    # -------- Counting --------
    def clear_results(self):
        for item in self.tree.get_children():
            self.tree.delete(item)
        self._results = []
        self._repetition = None
        self.export_btn.configure(state="disabled")
        self.export_md_btn.configure(state="disabled")
        self.copy_btn.configure(state="disabled")
        self.total_files_var.set("Files: 0")
        self.total_words_var.set("Total words: 0")
        self.total_unique_words_var.set("")
        self.total_rep_words_var.set("")
        self.total_billable_var.set("Unique units: 0")
        self.total_rep_billable_var.set("Rep. units: 0")
        self.total_amount_var.set("Total: 0.00")
        self.status_var.set(self._dependency_status())

    def run_count(self):
        s = self._get_settings()
        # if PDF got disabled, filter list now
        self._file_list = filter_supported(self._file_list, include_pdfs=s.pdf_include)

        if not self._file_list:
            folder = self.folder_var.get().strip()
            if folder and os.path.isdir(folder):
                self._file_list = iter_files(folder, s.include_subfolders, include_pdfs=s.pdf_include)

        if not self._file_list:
            self.clear_results()
            self.status_var.set("No supported files selected/found.")
            return

        # Clear UI for fresh results
        for item in self.tree.get_children():
            self.tree.delete(item)
        self._results = []
        self.export_btn.configure(state="disabled")
        self.export_md_btn.configure(state="disabled")
        self.copy_btn.configure(state="disabled")

        self.run_btn.configure(state="disabled")
        self.status_var.set("Counting…")

        worker = threading.Thread(target=self._worker, args=(list(self._file_list), s), daemon=True)
        worker.start()

    def _worker(self, paths: List[str], s: Settings):
        self._queue.put(("meta", len(paths)))

        for idx, p in enumerate(paths, start=1):
            result = extract_text_by_type(p, s)
            text, note, segments = result.text, result.note, result.segments
            if note and not text:
                # Error: no text could be extracted
                metrics = FileMetrics(p, 0, 0, 0, 0, 0, 0, 0.0, note)
            else:
                w, c, cns, n, sent, para, pages = compute_metrics(text, s)
                metrics = FileMetrics(p, w, c, cns, n, sent, para, pages, note or "", text, segments)
            self._queue.put(("result", metrics, idx, len(paths)))

        self._queue.put(("done",))

    def _poll_queue(self):
        try:
            while True:
                msg = self._queue.get_nowait()
                kind = msg[0]

                if kind == "meta":
                    n = msg[1]
                    self.total_files_var.set(f"Files: {n}")
                    self.total_words_var.set("Total words: 0")

                elif kind == "result":
                    m: FileMetrics = msg[1]
                    idx, n = msg[2], msg[3]
                    self._results.append(m)

                    # % numbers relative to words (FineCount-ish)
                    nums_pct = (m.numbers / m.words * 100.0) if m.words else 0.0

                    self.tree.insert(
                        "", "end",
                        values=(
                            m.words,
                            m.chars,
                            m.chars_nospace,
                            m.numbers,
                            f"{nums_pct:.2f}%",
                            m.sentences,
                            m.paragraphs,
                            f"{m.pages_est:.2f}",
                            m.note,
                            m.filepath
                        )
                    )

                    total_words = sum(r.words for r in self._results)
                    self.total_words_var.set(f"Total words: {total_words}")
                    self.status_var.set(f"Processed {idx}/{n}")

                    self.update_billing()

                elif kind == "done":
                    # Run cross-document repetition analysis
                    self._repetition = analyze_repetitions(self._results)
                    self._update_repetition_columns()

                    self.status_var.set(f"Done. {len(self._results)} file(s).")
                    self.run_btn.configure(state="normal")
                    has = "normal" if self._results else "disabled"
                    self.export_btn.configure(state=has)
                    self.export_md_btn.configure(state=has)
                    self.copy_btn.configure(state=has)
                    self.update_billing()

        except queue.Empty:
            pass

        self.after(100, self._poll_queue)

    def _update_repetition_columns(self):
        """Retroactively fill repetition columns after analysis completes."""
        if not self._repetition:
            return
        children = self.tree.get_children()
        for i, item_id in enumerate(children):
            if i >= len(self._results):
                break
            fp = self._results[i].filepath
            info = self._repetition.per_file.get(fp)
            if info:
                self.tree.set(item_id, "uniq_w", str(info.unique_words))
                self.tree.set(item_id, "rep_w", str(info.repeated_words))
                total_w = info.unique_words + info.repeated_words
                pct = (info.repeated_words / total_w * 100.0) if total_w > 0 else 0.0
                self.tree.set(item_id, "rep_pct", f"{pct:.1f}%")

        rep = self._repetition
        self.total_unique_words_var.set(f"Unique words: {rep.corpus_unique_words}")
        self.total_rep_words_var.set(f"Repeated words: {rep.corpus_repeated_words}")

    # -------- Billing --------
    def update_billing(self):
        if not self._results:
            self.total_billable_var.set("Unique units: 0")
            self.total_rep_billable_var.set("Rep. units: 0")
            self.total_amount_var.set("Total: 0.00")
            return

        bill_by = self.bill_by_var.get()
        rate = float(self.rate_var.get() or 0.0)
        rep_rate = float(self.rep_rate_var.get() or 0.0)
        tax = float(self.tax_var.get() or 0.0)
        discount = float(self.discount_var.get() or 0.0)
        currency = (self.currency_var.get() or "").strip() or "GBP"

        if self._repetition:
            rep = self._repetition
            if bill_by == "Words":
                unique_units = rep.corpus_unique_words
                rep_units = rep.corpus_repeated_words
            elif bill_by == "Characters":
                unique_units = rep.corpus_unique_chars
                rep_units = rep.corpus_repeated_chars
            else:  # Pages (est.)
                wpp = int(self.words_per_page_var.get() or 330)
                if wpp > 0:
                    unique_units = rep.corpus_unique_words / wpp
                    rep_units = rep.corpus_repeated_words / wpp
                else:
                    unique_units = 0
                    rep_units = 0
        else:
            # Before repetition analysis completes (during incremental loading)
            if bill_by == "Words":
                unique_units = sum(r.words for r in self._results)
            elif bill_by == "Characters":
                unique_units = sum(r.chars for r in self._results)
            else:
                unique_units = sum(r.pages_est for r in self._results)
            rep_units = 0

        subtotal = (unique_units * rate) + (rep_units * rep_rate)
        subtotal_after_discount = subtotal * (1.0 - (discount / 100.0))
        total = subtotal_after_discount * (1.0 + (tax / 100.0))

        fmt = lambda v: f"{v:.2f}" if isinstance(v, float) else str(v)
        self.total_billable_var.set(f"Unique units: {fmt(unique_units)}")
        self.total_rep_billable_var.set(f"Rep. units: {fmt(rep_units)}")
        self.total_amount_var.set(f"Total: {currency} {total:.2f}")

    # -------- Export --------
    def _format_clipboard_report(self) -> str:
        if not self._results:
            return ""

        def fit(text: str, width: int) -> str:
            t = str(text or "")
            if len(t) <= width:
                return t
            if width <= 1:
                return t[:width]
            return t[:width - 1] + "…"

        rows = []
        for r in self._results:
            nums_pct = (r.numbers / r.words * 100.0) if r.words else 0.0
            info = self._repetition.per_file.get(r.filepath) if self._repetition else None
            rep_pct = ""
            if info:
                total_w = info.unique_words + info.repeated_words
                rep_pct = f"{info.repeated_words / total_w * 100:.1f}%" if total_w > 0 else "0.0%"
            rows.append({
                "file": os.path.basename(r.filepath),
                "words": str(r.words),
                "uniq": str(info.unique_words) if info else "",
                "rep": str(info.repeated_words) if info else "",
                "rep_pct": rep_pct,
                "chars": str(r.chars),
                "chars_ns": str(r.chars_nospace),
                "nums": str(r.numbers),
                "nums_pct": f"{nums_pct:.2f}%",
                "sent": str(r.sentences),
                "para": str(r.paragraphs),
                "pages": f"{r.pages_est:.2f}",
                "note": r.note,
            })

        file_width = max(len("File"), *(len(row["file"]) for row in rows))
        note_width = 28

        cols = [
            ("file", "File", file_width, "l"),
            ("words", "Words", 8, "r"),
            ("uniq", "Unique", 8, "r"),
            ("rep", "Rep.", 8, "r"),
            ("rep_pct", "%Rep", 6, "r"),
            ("chars", "Chars", 8, "r"),
            ("chars_ns", "NoSp", 8, "r"),
            ("nums", "Nums", 6, "r"),
            ("nums_pct", "%Nums", 7, "r"),
            ("sent", "Sent", 6, "r"),
            ("para", "Para", 6, "r"),
            ("pages", "Pages", 7, "r"),
            ("note", "Note", note_width, "l"),
        ]

        def align(value: str, width: int, how: str) -> str:
            v = fit(value, width)
            return v.rjust(width) if how == "r" else v.ljust(width)

        header = " | ".join(align(title, width, "l") for _, title, width, _ in cols)
        sep = "-+-".join("-" * width for _, _, width, _ in cols)
        body = []
        for row in rows:
            line = " | ".join(align(row[key], width, how) for key, _, width, how in cols)
            body.append(line)

        total_words = sum(r.words for r in self._results)
        total_chars = sum(r.chars for r in self._results)
        total_pages = sum(r.pages_est for r in self._results)

        bill_by = self.bill_by_var.get()
        rate = float(self.rate_var.get() or 0.0)
        rep_rate = float(self.rep_rate_var.get() or 0.0)
        tax = float(self.tax_var.get() or 0.0)
        discount = float(self.discount_var.get() or 0.0)
        currency = (self.currency_var.get() or "").strip() or "GBP"

        if self._repetition:
            rep = self._repetition
            if bill_by == "Words":
                unique_units, rep_units = rep.corpus_unique_words, rep.corpus_repeated_words
            elif bill_by == "Characters":
                unique_units, rep_units = rep.corpus_unique_chars, rep.corpus_repeated_chars
            else:
                wpp = int(self.words_per_page_var.get() or 330)
                unique_units = rep.corpus_unique_words / wpp if wpp > 0 else 0
                rep_units = rep.corpus_repeated_words / wpp if wpp > 0 else 0
        else:
            if bill_by == "Words":
                unique_units = sum(r.words for r in self._results)
            elif bill_by == "Characters":
                unique_units = sum(r.chars for r in self._results)
            else:
                unique_units = sum(r.pages_est for r in self._results)
            rep_units = 0

        subtotal = (unique_units * rate) + (rep_units * rep_rate)
        subtotal_after_discount = subtotal * (1.0 - (discount / 100.0))
        total = subtotal_after_discount * (1.0 + (tax / 100.0))

        lines = [
            f"{APP_NAME} v{APP_VERSION} - Count Report",
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            "",
            header,
            sep,
            *body,
            "",
            f"Files: {len(self._results)}",
            f"Total words: {total_words}",
            f"Total chars: {total_chars}",
            f"Total pages (est.): {total_pages:.2f}",
        ]

        if self._repetition:
            rep = self._repetition
            lines.extend([
                "",
                f"Repetition analysis:",
                f"  Unique words: {rep.corpus_unique_words}  ({rep.corpus_unique_segments} segments)",
                f"  Repeated words: {rep.corpus_repeated_words}  ({rep.corpus_repeated_segments} segments)",
                f"  Distinct segments: {rep.corpus_unique_segments}  Total segments: {rep.corpus_total_segments}",
            ])

        unique_label = f"Unique units: {unique_units:.2f}" if isinstance(unique_units, float) else f"Unique units: {unique_units}"
        rep_label = f"Rep. units: {rep_units:.2f}" if isinstance(rep_units, float) else f"Rep. units: {rep_units}"
        lines.extend([
            "",
            f"Billing: {bill_by} | Rate {currency} {rate:.4f} | Rep. rate {currency} {rep_rate:.4f} | Discount {discount:.2f}% | Tax {tax:.2f}%",
            unique_label,
            rep_label,
            f"Total amount: {currency} {total:.2f}",
            "",
            "=" * 60,
            "DOCUMENT CONTENTS",
            "=" * 60,
        ])

        for r in self._results:
            fname = os.path.basename(r.filepath)
            lines.append("")
            lines.append(f"--- {fname} ---")
            lines.append("")
            if r.text and r.text.strip():
                lines.append(r.text.strip())
            else:
                lines.append(f"[{r.note or 'No text extracted'}]")

        return "\n".join(lines)

    def copy_report(self):
        if not self._results:
            return
        try:
            report = self._format_clipboard_report()
            self.clipboard_clear()
            self.clipboard_append(report)
            self.update()
            messagebox.showinfo("Copied", "Formatted report copied to clipboard.\nTip: use a fixed-width font in Gmail for best alignment.")
        except Exception as e:
            messagebox.showerror("Copy failed", str(e))

    def export_csv(self):
        if not self._results:
            return
        out = filedialog.asksaveasfilename(
            title="Save CSV",
            defaultextension=".csv",
            filetypes=[("CSV files", "*.csv")],
        )
        if not out:
            return

        try:
            with open(out, "w", newline="", encoding="utf-8") as f:
                w = csv.writer(f)
                w.writerow([
                    "file", "words", "unique_words", "repeated_words", "rep_pct",
                    "chars", "chars_no_spaces", "numbers", "sentences",
                    "paragraphs", "pages_est", "note"
                ])
                for r in self._results:
                    info = self._repetition.per_file.get(r.filepath) if self._repetition else None
                    uw = info.unique_words if info else ""
                    rw = info.repeated_words if info else ""
                    total_w = (info.unique_words + info.repeated_words) if info else 0
                    rp = f"{info.repeated_words / total_w * 100:.1f}%" if info and total_w > 0 else ""
                    w.writerow([r.filepath, r.words, uw, rw, rp, r.chars, r.chars_nospace, r.numbers,
                                r.sentences, r.paragraphs, f"{r.pages_est:.2f}", r.note])

                w.writerow([])
                total_words = sum(r.words for r in self._results)
                total_chars = sum(r.chars for r in self._results)
                total_pages = sum(r.pages_est for r in self._results)
                w.writerow(["TOTALS", total_words, "", "", "", total_chars, "", "", "", "", f"{total_pages:.2f}", ""])

                if self._repetition:
                    rep = self._repetition
                    w.writerow([])
                    w.writerow(["REPETITION SUMMARY"])
                    w.writerow(["Unique words", rep.corpus_unique_words])
                    w.writerow(["Repeated words", rep.corpus_repeated_words])
                    w.writerow(["Distinct segments", rep.corpus_unique_segments])
                    w.writerow(["Total segments", rep.corpus_total_segments])

            messagebox.showinfo("Exported", f"Saved:\n{out}")
        except Exception as e:
            messagebox.showerror("Export failed", str(e))

    def _format_markdown_report(self) -> str:
        if not self._results:
            return ""

        total_words = sum(r.words for r in self._results)
        total_chars = sum(r.chars for r in self._results)
        total_pages = sum(r.pages_est for r in self._results)

        bill_by = self.bill_by_var.get()
        rate = float(self.rate_var.get() or 0.0)
        rep_rate = float(self.rep_rate_var.get() or 0.0)
        tax = float(self.tax_var.get() or 0.0)
        discount = float(self.discount_var.get() or 0.0)
        currency = (self.currency_var.get() or "").strip() or "GBP"

        if self._repetition:
            rep = self._repetition
            if bill_by == "Words":
                unique_units, rep_units = rep.corpus_unique_words, rep.corpus_repeated_words
            elif bill_by == "Characters":
                unique_units, rep_units = rep.corpus_unique_chars, rep.corpus_repeated_chars
            else:
                wpp = int(self.words_per_page_var.get() or 330)
                unique_units = rep.corpus_unique_words / wpp if wpp > 0 else 0
                rep_units = rep.corpus_repeated_words / wpp if wpp > 0 else 0
        else:
            if bill_by == "Words":
                unique_units = sum(r.words for r in self._results)
            elif bill_by == "Characters":
                unique_units = sum(r.chars for r in self._results)
            else:
                unique_units = sum(r.pages_est for r in self._results)
            rep_units = 0

        subtotal = (unique_units * rate) + (rep_units * rep_rate)
        subtotal_after_discount = subtotal * (1.0 - (discount / 100.0))
        total = subtotal_after_discount * (1.0 + (tax / 100.0))

        lines = [
            f"# {APP_NAME} v{APP_VERSION} - Count Report",
            f"",
            f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            f"",
            f"## Summary",
            f"",
            f"| File | Words | Unique | Repeated | % Rep | Chars | Chars (no sp) | Numbers | % nums | Sent. | Para. | Pages (est.) | Note |",
            f"|------|------:|-------:|---------:|------:|------:|--------------:|--------:|-------:|------:|------:|-------------:|------|",
        ]

        for r in self._results:
            nums_pct = (r.numbers / r.words * 100.0) if r.words else 0.0
            fname = os.path.basename(r.filepath)
            info = self._repetition.per_file.get(r.filepath) if self._repetition else None
            uw = str(info.unique_words) if info else ""
            rw = str(info.repeated_words) if info else ""
            tw = (info.unique_words + info.repeated_words) if info else 0
            rp = f"{info.repeated_words / tw * 100:.1f}%" if info and tw > 0 else ""
            lines.append(
                f"| {fname} | {r.words} | {uw} | {rw} | {rp} | {r.chars} | {r.chars_nospace} "
                f"| {r.numbers} | {nums_pct:.2f}% | {r.sentences} | {r.paragraphs} "
                f"| {r.pages_est:.2f} | {r.note} |"
            )

        lines.extend([
            f"",
            f"**Files:** {len(self._results)}  ",
            f"**Total words:** {total_words}  ",
            f"**Total chars:** {total_chars}  ",
            f"**Total pages (est.):** {total_pages:.2f}",
        ])

        if self._repetition:
            rep = self._repetition
            lines.extend([
                f"",
                f"### Repetition Analysis",
                f"",
                f"- Unique words: {rep.corpus_unique_words} ({rep.corpus_unique_segments} segments)",
                f"- Repeated words: {rep.corpus_repeated_words} ({rep.corpus_repeated_segments} segments)",
                f"- Distinct segments: {rep.corpus_unique_segments} | Total segments: {rep.corpus_total_segments}",
            ])

        unique_label = f"{unique_units:.2f}" if isinstance(unique_units, float) else str(unique_units)
        rep_label = f"{rep_units:.2f}" if isinstance(rep_units, float) else str(rep_units)
        lines.extend([
            f"",
            f"### Billing",
            f"",
            f"- Bill by: {bill_by} | Rate: {currency} {rate:.4f} | Rep. rate: {currency} {rep_rate:.4f} | Discount: {discount:.2f}% | Tax: {tax:.2f}%",
            f"- Unique units: {unique_label} | Rep. units: {rep_label}",
            f"- **Total amount: {currency} {total:.2f}**",
            f"",
            f"---",
            f"",
            f"## Document Contents",
            f"",
        ])

        for r in self._results:
            fname = os.path.basename(r.filepath)
            lines.append(f"### {fname}")
            lines.append(f"")
            if r.text and r.text.strip():
                lines.append(r.text.strip())
            else:
                lines.append(f"*{r.note or 'No text extracted'}*")
            lines.append(f"")

        return "\n".join(lines)

    def export_md(self):
        if not self._results:
            return
        out = filedialog.asksaveasfilename(
            title="Save Markdown report",
            defaultextension=".md",
            filetypes=[("Markdown files", "*.md"), ("All files", "*.*")],
        )
        if not out:
            return

        try:
            report = self._format_markdown_report()
            with open(out, "w", encoding="utf-8") as f:
                f.write(report)
            messagebox.showinfo("Exported", f"Saved:\n{out}")
        except Exception as e:
            messagebox.showerror("Export failed", str(e))


if __name__ == "__main__":
    App().mainloop()